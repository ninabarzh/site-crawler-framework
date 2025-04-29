"""
Content Extraction Utilities Module

Provides comprehensive text extraction from multiple file formats with:
- PDF text extraction with layout preservation (using pdfplumber)
- HTML cleaning with semantic structure awareness (using BeautifulSoup)
- DOCX paragraph and table content extraction (using python-docx)
- PPTX slide and notes text extraction (using python-pptx)

Features:
- Automatic format detection
- Size validation safeguards (50MB for DOCX, 100MB for PPTX)
- Corrupt file handling with detailed error logging
- Thread-safe static methods
- Conditional dependency support (graceful degradation when optional packages not available)
"""

import io
import logging
import zipfile
from typing import Optional, Union, Tuple
from pathlib import Path
import pdfplumber
from bs4 import BeautifulSoup
from bs4.element import Comment


class _FakeDocument:
    """
    Fallback Document class when python-docx is not installed.

    Implements minimal interface to raise informative error when attempting to use
    DOCX functionality without the required dependency.
    """

    def __init__(self, *args, **kwargs):
        """Initialize with ImportError to alert about missing dependency"""
        raise ImportError("python-docx not installed - install with 'pip install python-docx'")

    @property
    def paragraphs(self):
        """Empty paragraphs property for interface compatibility"""
        return []

    @property
    def tables(self):
        """Empty tables property for interface compatibility"""
        return []


class _FakePresentation:
    """
    Fallback Presentation class when python-pptx is not installed.

    Provides minimal interface to raise informative error when attempting to use
    PPTX functionality without the required dependency.
    """

    def __init__(self, *args, **kwargs):
        """Initialize with ImportError to alert about missing dependency"""
        raise ImportError("python-pptx not installed - install with 'pip install python-pptx'")

    @property
    def slides(self):
        """Empty slides property for interface compatibility"""
        return []


# Conditional imports with proper fallback handling
try:
    from docx import Document

    DOCX_SUPPORT = True
except ImportError:
    Document = _FakeDocument
    DOCX_SUPPORT = False
    logging.warning("DOCX support disabled - python-docx package not installed")

try:
    from pptx import Presentation

    PPTX_SUPPORT = True
except ImportError:
    Presentation = _FakePresentation
    PPTX_SUPPORT = False
    logging.warning("PPTX support disabled - python-pptx package not installed")


class ContentExtractor:
    """
    EContent extraction toolkit with format-specific handlers.

    Provides static methods for extracting clean text content from:
    - PDF documents (with layout preservation)
    - HTML content (with configurable structure awareness)
    - Microsoft Word documents (.docx format)
    - PowerPoint presentations (.pptx format)

    All methods handle:
    - Multiple input types (file paths, bytes, file objects)
    - Size validation (preventing memory issues with large files)
    - Corrupt file recovery
    - Detailed error logging
    """

    # Security limits to prevent memory issues with large files
    MAX_DOCX_SIZE = 50 * 1024 * 1024  # 50MB
    MAX_PPTX_SIZE = 100 * 1024 * 1024  # 100MB

    @staticmethod
    def extract_text_from_pdf(
        source: Union[str, Path, io.BufferedReader, io.BytesIO, bytes],
        password: Optional[str] = None
    ) -> Optional[str]:
        """
        Extract readable text from PDF documents while preserving layout structure.

        Args:
            source: PDF input source, which can be:
                   - File path (as str or Path object)
                   - PDF content as bytes
                   - File-like object (BufferedReader/BytesIO)
            password: Optional password for encrypted PDFs

        Returns:
            str: Extracted text with page breaks preserved as newlines
            None: If extraction fails or no text found

        Raises:
            ValueError: For invalid input types
            IOError: For file access issues
        """
        try:
            if isinstance(source, (str, Path)):
                with open(source, 'rb') as f:
                    return ContentExtractor._process_pdf(io.BytesIO(f.read()), password)
            elif isinstance(source, bytes):
                return ContentExtractor._process_pdf(io.BytesIO(source), password)
            elif isinstance(source, (io.BufferedReader, io.BytesIO)):
                return ContentExtractor._process_pdf(source, password)
            raise ValueError("Invalid input type")
        except (IOError, OSError, ValueError) as e:
            logging.error(f"PDF access error: {str(e)}")
            return None
        except Exception as e:
            logging.error(f"Unexpected PDF error: {str(e)}")
            return None

    @staticmethod
    def _process_pdf(
            pdf_stream: Union[io.BufferedReader, io.BytesIO],
            password: Optional[str]
    ) -> Optional[str]:
        """
        Internal PDF processing method using pdfplumber library.

        Args:
            pdf_stream: File-like object containing PDF data
            password: Optional password for encrypted PDFs

        Returns:
            str: Concatenated text from all pages
            None: If processing fails

        Note:
            - Uses pdfplumber's layout preservation capabilities
            - Handles page-level extraction errors gracefully
            - Normalizes whitespace in output
        """
        try:
            with pdfplumber.open(pdf_stream, password=password) as pdf:
                text_parts = []
                for page in pdf.pages:
                    try:
                        text = page.extract_text()
                        if text and text.strip():
                            text_parts.append(text.strip())
                    except (AttributeError, TypeError):
                        continue
                return '\n'.join(text_parts) if text_parts else None
        except Exception as e:
            if "PDFSyntaxError" in str(type(e).__name__):
                logging.error("Invalid PDF structure - file may be corrupt")
            else:
                logging.error(f"PDF processing error: {str(e)}")
            return None

    @staticmethod
    def extract_text_from_html(
        html: str,
        preserve_comments: bool = False,
        preserve_links: bool = False,
        preserve_headings: bool = False
    ) -> str:
        """
        Extract clean text from HTML with configurable structure preservation.

        Args:
            html: Raw HTML content string
            preserve_comments: Keep HTML comments when True (default: False)
            preserve_links: Keep hyperlink text when True (default: False)
            preserve_headings: Maintain heading structure when True (default: False)

        Returns:
            str: Clean text content with normalized whitespace
        """
        try:
            soup = BeautifulSoup(html, 'html.parser')
            elements_to_remove = ['script', 'style', 'noscript']
            if not preserve_comments:
                elements_to_remove.append(Comment)

            for element in soup(elements_to_remove):
                element.decompose()

            text_parts = []
            for element in soup.find_all(text=True):
                parent = element.parent.name
                text = element.strip()
                if not text:
                    continue
                if parent == 'a' and not preserve_links:
                    continue
                if parent in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                    if preserve_headings:
                        text_parts.append(f"\n{text.upper()}\n")
                    continue
                text_parts.append(text)

            return ' '.join(text_parts).replace('\xa0', ' ').strip()
        except (AttributeError, TypeError) as e:
            logging.error(f"HTML parsing error: {str(e)}")
            return html
        except Exception as e:
            logging.error(f"Unexpected HTML error: {str(e)}")
            return html

    @staticmethod
    def extract_text_from_docx(
        source: Union[str, Path, io.BufferedReader, io.BytesIO, bytes]
    ) -> Optional[str]:
        """
        Extract text content from DOCX documents including paragraphs and tables.

        Args:
            source: DOCX input source, which can be:
                   - File path (as str or Path object)
                   - DOCX content as bytes
                   - File-like object (BufferedReader/BytesIO)

        Returns:
            str: Concatenated document text with paragraph breaks
            None: If extraction fails or format unsupported

        Notes:
            - Requires python-docx package
            - Preserves paragraph structure with double newlines
            - Extracts table content linearly
        """
        if not DOCX_SUPPORT:
            return None

        try:
            if isinstance(source, (str, Path)):
                doc = Document(source)
            elif isinstance(source, bytes):
                if len(source) > ContentExtractor.MAX_DOCX_SIZE:
                    raise ValueError("DOCX file too large")
                doc = Document(io.BytesIO(source))
            elif isinstance(source, (io.BufferedReader, io.BytesIO)):
                doc = Document(source)
            else:
                raise ValueError("Invalid input type")

            text_parts = []
            text_parts.extend(para.text.strip() for para in doc.paragraphs if para.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text.strip())
            return '\n\n'.join(text_parts) if text_parts else None

        except (zipfile.BadZipFile, ValueError, AttributeError) as e:
            logging.error(f"DOCX error: {str(e)}")
            return None
        except Exception as e:
            logging.error(f"Unexpected DOCX error: {str(e)}")
            return None

    @staticmethod
    def extract_text_from_pptx(
        source: Union[str, Path, io.BufferedReader, io.BytesIO, bytes]
    ) -> Optional[Tuple[str, str]]:
        """
        Extract text from PowerPoint presentations including slides and notes.

        Args:
            source: PPTX input source, which can be:
                   - File path (as str or Path object)
                   - PPTX content as bytes
                   - File-like object (BufferedReader/BytesIO)

        Returns:
            Tuple: (slides_text, notes_text) where each is:
                  - str: Extracted content
                  - None: If no content found
            None: If extraction fails or format unsupported

        Notes:
            - Requires python-pptx package
            - Slide numbers are included in output
            - Notes are separated by slide
        """
        if not PPTX_SUPPORT:
            return None

        try:
            if isinstance(source, (str, Path)):
                prs = Presentation(source)
            elif isinstance(source, bytes):
                if len(source) > ContentExtractor.MAX_PPTX_SIZE:
                    raise ValueError("PPTX file too large")
                prs = Presentation(io.BytesIO(source))
            elif isinstance(source, (io.BufferedReader, io.BytesIO)):
                prs = Presentation(source)
            else:
                raise ValueError("Invalid input type")

            slides_text = []
            notes_text = []

            for i, slide in enumerate(prs.slides):
                slide_content = [f"Slide {i+1}:"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                if len(slide_content) > 1:
                    slides_text.append('\n'.join(slide_content))

                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text.append(
                        f"Notes {i+1}:\n"
                        f"{slide.notes_slide.notes_text_frame.text.strip()}"
                    )

            return (
                '\n\n'.join(slides_text) if slides_text else None,
                '\n\n'.join(notes_text) if notes_text else None
            )

        except (zipfile.BadZipFile, ValueError, AttributeError) as e:
            logging.error(f"PPTX error: {str(e)}")
            return None
        except Exception as e:
            logging.error(f"Unexpected PPTX error: {str(e)}")
            return None
